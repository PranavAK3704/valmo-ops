"""
assessment-generator.py
=======================
Generates assessment questions from PPT using Groq API.
Runs in GitHub Actions after process_and_sync.py.

Flow:
1. Reads Training_Input sheet for new/updated processes
2. Downloads PPT from Google Drive
3. Extracts text from all slides
4. Calls Groq API to generate 5 assessment questions (3 MCQ + 2 subjective)
5. Writes questions to Assessment_Questions sheet in Google Sheets
6. Prevents duplicate questions for same process+version
"""

import os
import json
import csv
import tempfile
import requests
from io import StringIO

from pptx import Presentation
import gspread
from google.oauth2.service_account import Credentials


# ═══════════════════════════════════════════════════════════════
# CONFIG - from GitHub Secrets
# ═══════════════════════════════════════════════════════════════

SHEET_ID = os.environ["SHEET_ID"]
SHEET_INPUT_CSV_URL = os.environ["SHEET_INPUT_CSV_URL"]
SA_JSON = os.environ["GOOGLE_SERVICE_ACCOUNT_JSON"]
GROQ_API_KEY = os.environ["GROQ_API_KEY"]  # Add this to GitHub Secrets

INPUT_TAB_NAME = "Training_Input"
QUESTIONS_TAB_NAME = "Assessment_Questions"

GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "llama-3.3-70b-versatile"


# ═══════════════════════════════════════════════════════════════
# GOOGLE SHEETS CLIENT
# ═══════════════════════════════════════════════════════════════

def get_sheets_client():
    """Initialize Google Sheets client"""
    creds_dict = json.loads(SA_JSON)
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive.readonly"
    ]
    creds = Credentials.from_service_account_info(creds_dict, scopes=scopes)
    return gspread.authorize(creds)


# ═══════════════════════════════════════════════════════════════
# FETCH INPUT ROWS
# ═══════════════════════════════════════════════════════════════

def fetch_input_rows():
    """Fetch rows from Training_Input sheet"""
    print("📥 Fetching Training_Input rows...")
    response = requests.get(SHEET_INPUT_CSV_URL, timeout=15)
    response.raise_for_status()
    
    reader = csv.DictReader(StringIO(response.text))
    rows = list(reader)
    print(f"✅ Retrieved {len(rows)} rows")
    return rows


# ═══════════════════════════════════════════════════════════════
# PPT TEXT EXTRACTION
# ═══════════════════════════════════════════════════════════════

def extract_text_from_ppt(ppt_path):
    """Extract all text from PPT slides"""
    print(f"📄 Extracting text from PPT...")
    
    prs = Presentation(ppt_path)
    all_text = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        
        # Extract title
        if slide.shapes.title:
            slide_text.append(slide.shapes.title.text)
        
        # Extract all text from shapes
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_text.append(shape.text)
        
        if slide_text:
            all_text.append(f"Slide {slide_idx + 1}: " + " | ".join(slide_text))
    
    combined_text = "\n\n".join(all_text)
    print(f"✅ Extracted {len(combined_text)} characters from {len(prs.slides)} slides")
    
    return combined_text


# ═══════════════════════════════════════════════════════════════
# GROQ API - GENERATE QUESTIONS
# ═══════════════════════════════════════════════════════════════

def generate_questions_with_groq(process_name, ppt_text):
    """
    Use Groq API to generate assessment questions from PPT text.
    
    Returns:
        List of 5 questions (3 MCQ + 2 subjective)
    """
    print(f"🤖 Generating questions for: {process_name}")
    
    # Truncate text if too long (Groq has token limits)
    max_chars = 8000
    if len(ppt_text) > max_chars:
        ppt_text = ppt_text[:max_chars] + "\n\n[Text truncated due to length]"
    
    system_prompt = """You are an expert training assessment creator for logistics operations.

Generate 5 high-quality assessment questions from the provided training material.

Requirements:
1. Generate exactly 3 Multiple Choice Questions (MCQ)
2. Generate exactly 2 Subjective (essay) questions
3. MCQs should have 4 options each
4. Questions should test understanding, not just memorization
5. Subjective questions should require explanation and critical thinking
6. Questions should be clear, professional, and relevant to the process

Return ONLY valid JSON in this exact format (no markdown, no backticks):
{
  "questions": [
    {
      "type": "mcq",
      "question": "What is the first step in the process?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_index": 2,
      "points": 20
    },
    {
      "type": "subjective",
      "question": "Explain why this step is important.",
      "model_answer": "This step is important because...",
      "points": 25
    }
  ]
}"""

    user_prompt = f"""Process Name: {process_name}

Training Material:
{ppt_text}

Generate 5 assessment questions (3 MCQ + 2 subjective) based on this training material."""

    try:
        response = requests.post(
            GROQ_API_URL,
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": GROQ_MODEL,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.7,
                "max_tokens": 2000
            },
            timeout=30
        )
        
        if not response.ok:
            print(f"❌ Groq API error: {response.status_code}")
            print(f"Response: {response.text}")
            return None
        
        data = response.json()
        content = data["choices"][0]["message"]["content"]
        
        # Parse JSON response
        # Remove markdown code blocks if present
        content = content.strip()
        if content.startswith("```json"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()
        
        questions_data = json.loads(content)
        questions = questions_data["questions"]
        
        print(f"✅ Generated {len(questions)} questions")
        return questions
        
    except Exception as e:
        print(f"❌ Error generating questions: {e}")
        return None


# ═══════════════════════════════════════════════════════════════
# FALLBACK - DEFAULT QUESTIONS
# ═══════════════════════════════════════════════════════════════

def generate_default_questions(process_name):
    """Generate generic fallback questions if Groq fails"""
    print(f"⚠️  Using fallback questions for: {process_name}")
    
    return [
        {
            "type": "mcq",
            "question": f"What is the primary objective of the {process_name} process?",
            "options": [
                "Maximize speed",
                "Ensure accuracy and compliance",
                "Reduce costs",
                "All of the above"
            ],
            "correct_index": 3,
            "points": 20
        },
        {
            "type": "mcq",
            "question": "When should you escalate issues in this process?",
            "options": [
                "Never escalate",
                "Only for critical errors",
                "When uncertain about next steps",
                "After completing the process"
            ],
            "correct_index": 2,
            "points": 20
        },
        {
            "type": "mcq",
            "question": "What is the most important aspect of this process?",
            "options": [
                "Speed of completion",
                "Following documented procedures",
                "Customer satisfaction",
                "Cost reduction"
            ],
            "correct_index": 1,
            "points": 20
        },
        {
            "type": "subjective",
            "question": f"Describe the key steps involved in the {process_name} process and explain their importance.",
            "model_answer": f"The {process_name} process involves following established protocols, ensuring accuracy at each step, proper documentation, and timely completion. Each step is important to maintain quality standards and operational efficiency.",
            "points": 20
        },
        {
            "type": "subjective",
            "question": f"What are the potential consequences of errors in the {process_name} process?",
            "model_answer": "Errors can lead to operational delays, customer dissatisfaction, financial losses, compliance issues, and reduced efficiency. Proper execution ensures smooth operations and maintains service quality.",
            "points": 20
        }
    ]


# ═══════════════════════════════════════════════════════════════
# CHECK FOR EXISTING QUESTIONS
# ═══════════════════════════════════════════════════════════════

def check_existing_questions(client, process_name, version):
    """Check if questions already exist for this process+version"""
    try:
        sheet = client.open_by_key(SHEET_ID)
        
        try:
            worksheet = sheet.worksheet(QUESTIONS_TAB_NAME)
        except gspread.WorksheetNotFound:
            # Sheet doesn't exist yet, so no existing questions
            return False
        
        # Get all records
        records = worksheet.get_all_records()
        
        # Check if any match process_name and version
        for record in records:
            if (record.get("process_name") == process_name and 
                record.get("version") == version):
                return True
        
        return False
        
    except Exception as e:
        print(f"⚠️  Error checking existing questions: {e}")
        return False


# ═══════════════════════════════════════════════════════════════
# WRITE QUESTIONS TO GOOGLE SHEETS
# ═══════════════════════════════════════════════════════════════

def write_questions_to_sheet(client, process_name, version, questions):
    """Write generated questions to Assessment_Questions sheet"""
    print(f"📝 Writing {len(questions)} questions to sheet...")
    
    sheet = client.open_by_key(SHEET_ID)
    
    try:
        worksheet = sheet.worksheet(QUESTIONS_TAB_NAME)
    except gspread.WorksheetNotFound:
        # Create the sheet if it doesn't exist
        worksheet = sheet.add_worksheet(title=QUESTIONS_TAB_NAME, rows=500, cols=15)
        
        # Add headers
        headers = [
            "process_name",
            "version",
            "question_id",
            "type",
            "question",
            "option_a",
            "option_b", 
            "option_c",
            "option_d",
            "correct_index",
            "model_answer",
            "points",
            "generated_date"
        ]
        worksheet.update("A1", [headers])
    
    # Prepare rows to append
    rows = []
    for idx, q in enumerate(questions, start=1):
        question_id = f"{process_name.lower().replace(' ', '_')}_q{idx}"
        
        if q["type"] == "mcq":
            row = [
                process_name,
                version,
                question_id,
                "mcq",
                q["question"],
                q["options"][0],
                q["options"][1],
                q["options"][2],
                q["options"][3],
                q["correct_index"],
                "",  # No model_answer for MCQ
                q["points"],
                "2026-03-02"  # Current date
            ]
        else:  # subjective
            row = [
                process_name,
                version,
                question_id,
                "subjective",
                q["question"],
                "", "", "", "",  # No options for subjective
                "",  # No correct_index
                q["model_answer"],
                q["points"],
                "2026-03-02"
            ]
        
        rows.append(row)
    
    # Append all rows
    worksheet.append_rows(rows)
    print(f"✅ Written {len(rows)} questions to sheet")


# ═══════════════════════════════════════════════════════════════
# DOWNLOAD PPT FROM GOOGLE DRIVE
# ═══════════════════════════════════════════════════════════════

def extract_drive_file_id(url):
    """Extract file ID from Google Drive URL"""
    import re
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    match = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    raise ValueError(f"Could not extract file ID from: {url}")


def download_drive_file(drive_url, save_path):
    """Download PPT from Google Drive"""
    file_id = extract_drive_file_id(drive_url)
    download_url = f"https://drive.google.com/uc?export=download&id={file_id}"
    
    session = requests.Session()
    response = session.get(download_url, stream=True)
    
    # Handle Google Drive virus scan warning
    for key, value in response.cookies.items():
        if key.startswith('download_warning'):
            download_url = f"https://drive.google.com/uc?export=download&id={file_id}&confirm={value}"
            response = session.get(download_url, stream=True)
            break
    
    with open(save_path, "wb") as f:
        for chunk in response.iter_content(chunk_size=32768):
            if chunk:
                f.write(chunk)
    
    file_size = os.path.getsize(save_path)
    print(f"✅ Downloaded PPT ({file_size:,} bytes)")


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    print("\n" + "="*60)
    print("🤖 ASSESSMENT GENERATOR - Groq AI")
    print("="*60)
    
    # Fetch input rows
    rows = fetch_input_rows()
    
    if not rows:
        print("❌ No rows in Training_Input sheet")
        return
    
    # Initialize Google Sheets client
    client = get_sheets_client()
    
    generated_count = 0
    skipped_count = 0
    failed_count = 0
    
    for idx, row in enumerate(rows, start=1):
        process_name = row.get("process_name", "").strip()
        ppt_link = row.get("ppt_link", "").strip()
        version = row.get("version", "1.0").strip()
        
        if not ppt_link or not process_name:
            print(f"\n⚠️  Row {idx}: Missing data, skipping")
            skipped_count += 1
            continue
        
        print(f"\n{'='*60}")
        print(f"Processing: {process_name} (v{version})")
        print(f"{'='*60}")
        
        # Check if questions already exist
        if check_existing_questions(client, process_name, version):
            print(f"✅ Questions already exist for {process_name} v{version}, skipping")
            skipped_count += 1
            continue
        
        # Download PPT
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            temp_path = tmp.name
        
        try:
            download_drive_file(ppt_link, temp_path)
            
            # Extract text
            ppt_text = extract_text_from_ppt(temp_path)
            
            # Generate questions with Groq
            questions = generate_questions_with_groq(process_name, ppt_text)
            
            # Use fallback if Groq fails
            if not questions:
                questions = generate_default_questions(process_name)
            
            # Write to sheet
            write_questions_to_sheet(client, process_name, version, questions)
            
            generated_count += 1
            print(f"✅ Success: {process_name}")
            
        except Exception as e:
            print(f"❌ Error processing {process_name}: {e}")
            failed_count += 1
            
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)
    
    # Summary
    print(f"\n{'='*60}")
    print(f"✅ ASSESSMENT GENERATION COMPLETE")
    print(f"{'='*60}")
    print(f"Generated: {generated_count}")
    print(f"Skipped: {skipped_count}")
    print(f"Failed: {failed_count}")
    print(f"Total: {len(rows)}")


if __name__ == "__main__":
    main()