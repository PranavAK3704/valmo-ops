"""
Platform Detector
=================
Identifies whether a process is Log10 or external platform.
Only Log10 processes get step extraction.

Strategy:
- Scan all text in process slides
- Match against known keyword sets
- Return platform type
"""

from pptx import Presentation

# Log10 platform indicators
LOG10_KEYWORDS = [
    'log10', 'log 10', 'log-10',
    'trips', 'manifest', 'bagging', 'rto',
    'shipment', 'awb', 'hub',
    'inbound', 'outbound', 'inventory',
    'forward', 'reverse', 'pickup', 'delivery',
    'scan', 'barcode', 'courier',
    'loadshare'  # Company name often in Log10 training
]

# External platform indicators
EXTERNAL_PLATFORMS = {
    'euphoria': ['euphoria', 'order consumables', 'buy now', 'add to cart', 'valmo'],
    'ticketing': ['ticket', 'kapture', 'support portal', 'raise ticket'],
    'email': ['outlook', 'gmail', 'email', 'pre-alert', 'mail'],
    'excel': ['excel', 'spreadsheet', 'csv', 'worksheet']
}


def get_segment_text(segment, ppt_path):
    """
    Extracts all text from slides in a segment.
    
    Args:
        segment: Process segment dict with "Slides" list
        ppt_path: Path to PPT file
    
    Returns:
        Combined text from all slides in segment
    """
    prs = Presentation(ppt_path)
    all_text = []
    
    # Include title
    all_text.append(segment.get("Process_Name", ""))
    
    # Extract text from all slides in segment
    for slide_idx in segment.get("Slides", []):
        if slide_idx >= len(prs.slides):
            continue
            
        slide = prs.slides[slide_idx]
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                all_text.append(shape.text)
    
    return " ".join(all_text)


def detect_platform(segment, ppt_path):
    """
    Detects which platform a process belongs to.
    
    Args:
        segment: Process segment from ppt_segmenter
        ppt_path: Path to PPT file
    
    Returns:
        "log10" | "external" | "unknown"
    """
    # Get all text from this process
    text = get_segment_text(segment, ppt_path).lower()
    
    # Check for Log10 indicators
    log10_score = sum(1 for kw in LOG10_KEYWORDS if kw in text)
    
    # Check for external platform indicators
    external_score = 0
    detected_platform = None
    
    for platform, keywords in EXTERNAL_PLATFORMS.items():
        score = sum(1 for kw in keywords if kw in text)
        if score > external_score:
            external_score = score
            detected_platform = platform
    
    # Decision logic
    if log10_score >= 2:  # At least 2 Log10 keywords
        return "log10"
    elif external_score >= 2:  # At least 2 external keywords
        return detected_platform
    elif log10_score > 0:  # At least 1 Log10 keyword
        return "log10"
    else:
        return "unknown"


if __name__ == "__main__":
    # Test platform detection
    import sys
    import os
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from ppt_segmenter import segment_processes
    
    test_ppts = [
        "data/ppts/RTO Bagging.pptx",
        "data/ppts/LM_Captains_Order_Consumables.pptx"
    ]
    
    for ppt_path in test_ppts:
        if not os.path.exists(ppt_path):
            continue
            
        print(f"\n{'='*60}")
        print(f"Testing: {os.path.basename(ppt_path)}")
        print(f"{'='*60}")
        
        segments = segment_processes(ppt_path)
        
        for segment in segments:
            platform = detect_platform(segment, ppt_path)
            print(f"\nProcess: {segment['Process_Name']}")
            print(f"Platform: {platform}")