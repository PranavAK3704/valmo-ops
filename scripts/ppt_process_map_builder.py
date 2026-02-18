"""
PPT Process Map Builder - Log10 Focused
========================================
Extracts actionable process intelligence ONLY for Log10 processes.
External processes are logged separately for reference.

NOW POWERED BY: Google Sheets + Google Drive

Strategy:
1. Detect platform (Log10 vs External)
2. For Log10: Extract start tab + step sequence
3. For External: Skip step extraction, just register
4. Output separate files for each

Key Innovation:
- Platform-aware processing
- Only extracts steps where they matter (Log10)
- Graceful degradation if extraction fails
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os
import sys
import re
import requests
import csv
from io import StringIO

# Add scripts to path for imports
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from ppt_segmenter import segment_processes
from platform_detector import detect_platform

# Red box detection - adjusted for actual PPT colors
RED_THRESHOLD = (150, 0, 0)

# ═══════════════════════════════════════════════════════════════
# GOOGLE SHEETS + DRIVE INTEGRATION (NEW)
# ═══════════════════════════════════════════════════════════════

# Your published Google Sheet CSV URL
SHEET_CSV_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vTKDX4o1H_sZSJS_tUDo68N1SyjV3m3kbnkucjLe-4y1cUR3PBb2O49fbfNe2AQt-Oiuiu0Egj-wi_P/pub?gid=0&single=true&output=csv"


def fetch_sheet_rows():
    """
    Fetches training data from published Google Sheet.
    
    Returns:
        List of dicts with keys: process_name, ppt_link, video_link
    """
    print("🌐 Fetching training data from Google Sheets...")
    try:
        response = requests.get(SHEET_CSV_URL, timeout=10)
        response.raise_for_status()
        
        reader = csv.DictReader(StringIO(response.text))
        rows = list(reader)
        
        print(f"✅ Retrieved {len(rows)} rows from sheet")
        return rows
    except Exception as e:
        print(f"❌ Failed to fetch Google Sheet: {e}")
        return []


def extract_drive_file_id(url):
    """
    Extracts file ID from Google Drive share URL.
    
    Supports formats:
    - https://drive.google.com/file/d/FILE_ID/view?usp=sharing
    - https://drive.google.com/open?id=FILE_ID
    """
    # Pattern 1: /d/FILE_ID/
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    
    # Pattern 2: ?id=FILE_ID
    match = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    
    raise ValueError(f"Could not extract file ID from: {url}")


def download_drive_file(drive_url, save_path):
    """
    Downloads a file from Google Drive to local path.
    
    Args:
        drive_url: Google Drive share link
        save_path: Local path to save file
    """
    try:
        file_id = extract_drive_file_id(drive_url)
        download_url = f"https://drive.google.com/uc?export=download&id={file_id}"
        
        # First request
        session = requests.Session()
        response = session.get(download_url, stream=True)
        
        # Check if we need to confirm download (large files)
        for key, value in response.cookies.items():
            if key.startswith('download_warning'):
                # Get confirmation token
                download_url = f"https://drive.google.com/uc?export=download&id={file_id}&confirm={value}"
                response = session.get(download_url, stream=True)
                break
        
        # Save file
        with open(save_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=32768):
                if chunk:
                    f.write(chunk)
        
        # Verify it's a valid PPT file (not HTML error page)
        file_size = os.path.getsize(save_path)
        if file_size < 1000:
            # Suspiciously small - might be error page
            with open(save_path, 'rb') as f:
                header = f.read(100)
                if b'<!DOCTYPE' in header or b'<html' in header:
                    raise ValueError("Downloaded HTML instead of PPT - check Drive permissions")
        
        print(f"⬇️  Downloaded PPT → {save_path} ({file_size:,} bytes)")
        
    except Exception as e:
        print(f"❌ Download failed: {e}")
        raise


# ═══════════════════════════════════════════════════════════════
# ORIGINAL EXTRACTION LOGIC (UNCHANGED)
# ═══════════════════════════════════════════════════════════════

# ─── URL Module Resolver ───
# Reads log10_tab_url_map.csv, extracts the module segment from each URL,
# and builds a lookup so start_tab → url_module is automatic.

TAB_URL_MAP_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "data", "log10_tab_url_map.csv")

def _extract_module_from_url(url: str):
    """
    Pulls the module segment out of a Log10 URL.
    https://log10-atlas.loadshare.net/appv2/rto/dashboard/waybill
                                            ^^^  ← this
    """
    try:
        path = url.split("?")[0]
        parts = [p for p in path.split("/") if p]
        appv2_idx = parts.index("operations")
        return parts[appv2_idx + 1]
    except (ValueError, IndexError):
        return None


def load_tab_url_map(csv_path=TAB_URL_MAP_PATH):
    """
    Returns dict { normalised_tab_name: url_module }
    e.g. { "rto": "rto", "sc-ops dashboard": "sc-ops", "tracking": "tracking" }
    """
    import csv

    mapping = {}

    if not os.path.exists(csv_path):
        print(f"⚠️  Tab-URL map not found at {csv_path} — url_module will be None")
        return mapping

    with open(csv_path, newline="", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            tab    = row.get("tab", "").strip().lower()
            url    = row.get("url", "").strip()
            module = _extract_module_from_url(url)
            if tab and module:
                mapping[tab] = module

    print(f"📋 Loaded {len(mapping)} tab→module mappings from {csv_path}")
    return mapping


def resolve_url_module(start_tab, tab_map):
    """
    Given a start_tab (e.g. "RTO", "Dashboard") returns the url_module or None.

    Match order:
      1. Exact match on lowercased start_tab        → "rto"  matches key "rto"
      2. start_tab is a substring of a map key      → "dashboard" inside "sc-ops dashboard"
      3. A map key is a substring of start_tab      → fallback
    """
    needle = start_tab.strip().lower()

    # 1. exact
    if needle in tab_map:
        return tab_map[needle]

    # 2. needle inside a key — pick longest key (most specific)
    candidates = [(k, v) for k, v in tab_map.items() if needle in k]
    if candidates:
        candidates.sort(key=lambda x: len(x[0]), reverse=True)
        return candidates[0][1]

    # 3. key inside needle
    candidates = [(k, v) for k, v in tab_map.items() if k in needle]
    if candidates:
        candidates.sort(key=lambda x: len(x[0]), reverse=True)
        return candidates[0][1]

    return None


def clean_process_name(title_text):
    """
    Cleans process name by removing newlines and generic suffixes.
    
    Example:
    "Order Consumables\nTraining Material" → "Order Consumables"
    """
    # Remove newlines
    name = title_text.replace('\n', ' ').replace('\r', ' ')
    
    # Normalize whitespace
    name = ' '.join(name.split())
    
    # Remove generic suffixes
    suffixes = ['Training Material', 'Training', 'Process', 'SOP', 'Procedure']
    for suffix in suffixes:
        if name.endswith(suffix):
            name = name[:-len(suffix)].strip()
    
    return name


def clean_tab_name(raw_tab):
    """
    Cleans tab name by removing step numbers, action verbs, newlines.
    
    Examples:
    "3. Click on\nLogin" → "Login"
    "Go to RTO(1)" → "RTO"
    """
    if not raw_tab:
        return None
    
    # Remove newlines
    tab = raw_tab.replace('\n', ' ').replace('\r', ' ')
    
    # Normalize whitespace
    tab = ' '.join(tab.split())
    
    # Remove step numbers at start
    tab = re.sub(r'^\d+\.\s*', '', tab)
    
    # Remove action verbs
    action_verbs = [
        'Go to', 'go to',
        'Click on', 'click on',
        'Then click', 'then click',
        'Click', 'click',
        'Select', 'select',
        'Open', 'open',
        'Choose', 'choose',
        'Navigate to', 'navigate to'
    ]
    
    for verb in action_verbs:
        if tab.startswith(verb):
            tab = tab[len(verb):].strip()
            break
    
    # Remove trailing (N) patterns
    tab = re.sub(r'\s*\(\d+\)\s*$', '', tab)
    
    return tab.strip()


def is_red_outlined_shape(shape):
    """Detects shapes with red outlines."""
    try:
        if not hasattr(shape, 'line'):
            return False
        
        if not shape.line.color.rgb:
            return False
        
        r, g, b = shape.line.color.rgb
        return r >= RED_THRESHOLD[0] and g < 100 and b < 100
    except:
        return False


def extract_text_from_shape(shape):
    """Safely extracts text from a shape."""
    try:
        if hasattr(shape, 'text') and shape.text:
            return shape.text.strip()
    except:
        pass
    return None


def find_nearest_text(slide, ref_shape, max_distance=2000000):
    """Finds text nearest to a red-outlined shape."""
    min_dist = float("inf")
    best_text = None
    
    for shape in slide.shapes:
        if shape == ref_shape:
            continue
        
        text = extract_text_from_shape(shape)
        if not text or len(text) > 100:
            continue
        
        dx = abs(shape.left - ref_shape.left)
        dy = abs(shape.top - ref_shape.top)
        dist = dx + dy
        
        if dist < min_dist and dist < max_distance:
            min_dist = dist
            best_text = text
    
    return best_text


def extract_instruction_text(slide):
    """
    Extracts explicit instruction text from slide.
    Looks for sentences with instruction markers.
    """
    for shape in slide.shapes:
        text = extract_text_from_shape(shape)
        if text and len(text) > 50:
            instruction_markers = ['go to', 'click on', 'click', 'then', 'select']
            if any(marker in text.lower() for marker in instruction_markers):
                return text
    return None


def parse_instruction_sequence(instruction_text):
    """
    Parses instruction text to extract Log10 tab sequence.
    
    Pattern: "Go to RTO(1). Click on RTO Manifest(2)..."
    Returns: ['RTO', 'RTO Manifest', 'Create Manifest']
    """
    if not instruction_text:
        return []
    
    # Find patterns like "RTO(1)", "RTO Manifest(2)"
    pattern = r'([A-Za-z][A-Za-z\s]+)\(\d+\)'
    matches = re.findall(pattern, instruction_text)
    
    clean_steps = []
    for match in matches:
        cleaned = clean_tab_name(match)
        if cleaned and cleaned not in clean_steps:
            clean_steps.append(cleaned)
    
    return clean_steps


def extract_log10_steps(segment, prs):
    """
    Extracts Log10 tab sequence from a process segment.
    
    Returns:
        (start_tab, steps_list) or (None, []) if extraction fails
    """
    steps = []
    start_tab = None
    instruction_sequence = []
    
    for slide_idx in segment["Slides"]:
        if slide_idx >= len(prs.slides):
            continue
        
        slide = prs.slides[slide_idx]
        
        # Method 1: Extract from instruction text (primary)
        instruction_text = extract_instruction_text(slide)
        if instruction_text:
            sequence = parse_instruction_sequence(instruction_text)
            if sequence:
                instruction_sequence.extend(sequence)
        
        # Method 2: Detect red boxes (fallback)
        for shape in slide.shapes:
            if is_red_outlined_shape(shape):
                text = find_nearest_text(slide, shape)
                if text:
                    cleaned = clean_tab_name(text)
                    if cleaned and cleaned not in steps:
                        steps.append(cleaned)
    
    # Prefer instruction sequence (more reliable)
    if instruction_sequence:
        start_tab = instruction_sequence[0]
        steps = instruction_sequence
    elif steps:
        start_tab = steps[0]
    
    return start_tab, steps


def build_process_map(ppt_path, demo_video_link):
    """
    Main Brain function - platform-aware processing.
    
    Returns:
        {
            "log10": [...],      # Processes for overlay
            "external": [...]    # Processes for reference
        }
    """
    prs = Presentation(ppt_path)
    segments = segment_processes(ppt_path)

    # Load tab→module map once
    tab_map = load_tab_url_map()
    
    if not segments:
        print("⚠️  No processes detected in PPT")
        return {"log10": [], "external": []}
    
    log10_processes = []
    external_processes = []
    
    for segment in segments:
        # Detect platform
        platform = detect_platform(segment, ppt_path)
        
        process_name = clean_process_name(segment["Process_Name"])
        
        print(f"\n🔍 Analyzing: {process_name}")
        print(f"   Platform: {platform}")
        
        if platform == "log10":
            # Extract Log10 tabs
            try:
                start_tab, steps = extract_log10_steps(segment, prs)
                
                if start_tab:
                    url_module = resolve_url_module(start_tab, tab_map)
                    log10_processes.append({
                        "process_name": process_name,
                        "platform": "log10",
                        "start_tab": start_tab,
                        "url_module": url_module,
                        "steps": steps,
                        "video_link": demo_video_link,
                        **({"needs_review": True} if not url_module else {})
                    })
                    print(f"   ✅ Start Tab: {start_tab}")
                    print(f"   ✅ URL Module: {url_module}")
                    print(f"   ✅ Steps: {len(steps)}")
                else:
                    # Couldn't extract start tab - still add but flag
                    log10_processes.append({
                        "process_name": process_name,
                        "platform": "log10",
                        "start_tab": "Dashboard",  # Safe default
                        "url_module": resolve_url_module("Dashboard", tab_map),
                        "steps": [],
                        "video_link": demo_video_link,
                        "needs_review": True
                    })
                    print(f"   ⚠️  Could not extract start tab - needs review")
            
            except Exception as e:
                print(f"   ❌ Error extracting steps: {e}")
                # Still add process with defaults
                log10_processes.append({
                    "process_name": process_name,
                    "platform": "log10",
                    "start_tab": "Dashboard",
                    "url_module": resolve_url_module("Dashboard", tab_map),
                    "steps": [],
                    "video_link": demo_video_link,
                    "needs_review": True
                })
        
        else:
            # External process - just register, no step extraction
            external_processes.append({
                "process_name": process_name,
                "platform": platform if platform != "unknown" else "external",
                "video_link": demo_video_link,
                "use_case": "training_only"
            })
            print(f"   ℹ️  External process - no step extraction")
    
    return {
        "log10": log10_processes,
        "external": external_processes
    }


BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(BASE_DIR, "data", "output")

def save_process_maps(results, output_dir=OUTPUT_DIR):
    """Saves separate JSON files for Log10 and external processes."""
    os.makedirs(output_dir, exist_ok=True)
    
    # Save Log10 processes (for overlay)
    log10_path = os.path.join(output_dir, "log10_processes.json")
    with open(log10_path, "w") as f:
        json.dump(results["log10"], f, indent=2)
    print(f"\n📄 Log10 processes saved to: {log10_path}")
    
    # Save external processes (for reference)
    if results["external"]:
        external_path = os.path.join(output_dir, "external_processes.json")
        with open(external_path, "w") as f:
            json.dump(results["external"], f, indent=2)
        print(f"📄 External processes saved to: {external_path}")
    
    # Also save legacy format for backward compatibility
    legacy = results["log10"] + results["external"]
    legacy_path = os.path.join(output_dir, "process_map.json")
    with open(legacy_path, "w") as f:
        json.dump(legacy, f, indent=2)
    print(f"📄 Legacy format saved to: {legacy_path}")


# ═══════════════════════════════════════════════════════════════
# MAIN ENTRY POINT (MODIFIED FOR GOOGLE SHEETS)
# ═══════════════════════════════════════════════════════════════

def build_all_process_maps():
    """
    NOW POWERED BY GOOGLE SHEETS:
    1. Fetch training data from published Google Sheet
    2. Download each PPT from Drive temporarily
    3. Run existing extraction logic
    4. Inject real video links from sheet
    5. Clean up temp files
    """
    
    # Fetch from Google Sheets
    rows = fetch_sheet_rows()
    
    if not rows:
        print("❌ No data retrieved from Google Sheet")
        return
    
    all_log10 = []
    all_external = []
    
    for idx, row in enumerate(rows, start=1):
        process_name = row.get("process_name", "").strip()
        ppt_link = row.get("ppt_link", "").strip()
        video_link = row.get("video_link", "").strip()
        
        if not ppt_link:
            print(f"\n⚠️  Row {idx}: Skipping (no PPT link)")
            continue
        
        if not process_name:
            process_name = f"Process_{idx}"
        
        # Temp file path
        temp_path = f"temp_ppt_{idx}.pptx"
        
        try:
            print(f"\n{'='*60}")
            print(f"Row {idx}: {process_name}")
            print(f"{'='*60}")
            
            # Download PPT from Drive
            download_drive_file(ppt_link, temp_path)
            
            # Process using existing brain logic
            results = build_process_map(temp_path, video_link)
            
            all_log10.extend(results["log10"])
            all_external.extend(results["external"])
            
        except Exception as e:
            print(f"❌ Error processing {process_name}: {e}")
            import traceback
            traceback.print_exc()
        
        finally:
            # Always clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print(f"🗑️  Cleaned up: {temp_path}")
    
    # Save results
    if all_log10 or all_external:
        save_process_maps({
            "log10": all_log10,
            "external": all_external
        })
        
        print(f"\n{'='*60}")
        print(f"✅ SUCCESS")
        print(f"{'='*60}")
        print(f"Log10 processes: {len(all_log10)}")
        print(f"External processes: {len(all_external)}")
        
        # Summary
        if all_log10:
            print(f"\n📊 Log10 Processes (for overlay):")
            for proc in all_log10:
                review_flag = " ⚠️ NEEDS REVIEW" if proc.get("needs_review") else ""
                print(f"  • {proc['process_name']}")
                print(f"    Start: {proc['start_tab']}{review_flag}")
                print(f"    Module: {proc.get('url_module', '⚠️ None')}")
                print(f"    Steps: {len(proc['steps'])}")
        
        if all_external:
            print(f"\n📚 External Processes (training only):")
            for proc in all_external:
                print(f"  • {proc['process_name']} ({proc['platform']})")
    else:
        print("\n⚠️  No processes could be extracted")


if __name__ == "__main__":
    build_all_process_maps()