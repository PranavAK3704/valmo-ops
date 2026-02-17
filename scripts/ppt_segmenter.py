"""
PPT Segmentation Script
=======================
Detects process boundaries using slide layout patterns.
More robust than color-based detection.

Strategy:
- Title slides use "Title Slide" or "Divider Slide" layouts
- Content slides use "Custom Layout" or other layouts
- First title slide = main process
- Subsequent title slides (if any) = sub-processes
"""

from pptx import Presentation


def is_title_slide(slide):
    """
    Detects if a slide is a title/divider slide using layout name.
    
    This is more reliable than color detection because:
    - Layout names are explicit
    - Works with any theme
    - Device-independent
    """
    layout_name = slide.slide_layout.name.lower()
    
    # Common title slide layout names
    title_keywords = ['title', 'divider', 'section']
    
    is_title = any(keyword in layout_name for keyword in title_keywords)
    
    # Filter out closing slides
    if is_title and slide.shapes.title:
        title_text = slide.shapes.title.text.lower()
        closing_keywords = ['thank', 'questions', 'q&a', 'end', 'conclusion']
        if any(keyword in title_text for keyword in closing_keywords):
            return False
    
    return is_title


def is_content_slide(slide):
    """
    Detects content slides (process steps with screenshots).
    """
    layout_name = slide.slide_layout.name.lower()
    
    # Content slides typically have "custom" or "content" in the name
    content_keywords = ['custom', 'content', 'blank']
    
    # Also check for images (screenshots)
    has_images = any(shape.shape_type == 13 for shape in slide.shapes)
    
    return any(keyword in layout_name for keyword in content_keywords) or has_images


def segment_processes(ppt_path):
    """
    Segments a PPT into logical processes based on slide layout.
    
    Returns:
        List of process segments, each containing:
        - Process_Name: Title text
        - Start_Slide: Index where process begins
        - Slides: List of slide indices in this process
    """
    prs = Presentation(ppt_path)
    
    segments = []
    current = None
    
    for idx, slide in enumerate(prs.slides):
        # Skip slides without titles
        if not slide.shapes.title:
            if current:
                current["Slides"].append(idx)
            continue
        
        title_text = slide.shapes.title.text.strip()
        
        # Check if this is a title slide
        if is_title_slide(slide):
            # Close previous segment
            if current:
                segments.append(current)
            
            # Start new segment
            current = {
                "Process_Name": title_text,
                "Start_Slide": idx,
                "Slides": []
            }
            
            print(f"🎯 Process detected: '{title_text}' (Slide {idx})")
        
        elif current:
            # Add to current process
            current["Slides"].append(idx)
    
    # Don't forget the last segment
    if current:
        segments.append(current)
    
    print(f"\n📊 Total processes: {len(segments)}")
    return segments


if __name__ == "__main__":
    # Test segmentation
    import sys
    import os
    
    test_ppt = "data/ppts/RTO Bagging.pptx"
    
    if os.path.exists(test_ppt):
        print(f"Testing segmentation on: {test_ppt}\n")
        segments = segment_processes(test_ppt)
        
        print("\nSegmentation results:")
        print("="*60)
        for seg in segments:
            print(f"\n{seg['Process_Name']}")
            print(f"  Starts at: Slide {seg['Start_Slide']}")
            print(f"  Contains slides: {seg['Slides']}")
    else:
        print(f"❌ Test file not found: {test_ppt}")