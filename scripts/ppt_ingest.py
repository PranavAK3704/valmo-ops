"""
PPT Ingestion Script
====================
Registers PPT files and their associated demo video links.
Creates/updates brain_manifest.csv with metadata.

Run: python scripts/ppt_ingest.py
"""

import pandas as pd
from pptx import Presentation
import os
from datetime import datetime, timezone

MANIFEST_PATH = "data/output/brain_manifest.csv"
PPT_DIR = "data/ppts"
DEMO_LINKS_FILE = "data/demos/demo_links.json"


def ingest_ppt(ppt_path, demo_video_link):
    """
    Registers a single PPT file with its video link.
    
    Args:
        ppt_path: Path to the .pptx file
        demo_video_link: Associated video URL or placeholder
    """
    try:
        prs = Presentation(ppt_path)
        
        # Extract title from first slide
        extracted_title = "Unknown Process"
        if prs.slides and prs.slides[0].shapes.title:
            extracted_title = prs.slides[0].shapes.title.text.strip()
        
        entry = {
            "Process_Title": extracted_title,
            "PPT_File": ppt_path,
            "Demo_Video_Link": demo_video_link,
            "Ingested_At": datetime.now(timezone.utc).isoformat(),
            "Slide_Count": len(prs.slides)
        }
        
        df = pd.DataFrame([entry])
        
        # Create output directory if needed
        os.makedirs(os.path.dirname(MANIFEST_PATH), exist_ok=True)
        
        # Append to manifest (create if doesn't exist)
        df.to_csv(
            MANIFEST_PATH,
            mode='a',
            header=not os.path.exists(MANIFEST_PATH),
            index=False
        )
        
        print(f"✅ Ingested: {extracted_title} ({len(prs.slides)} slides)")
        return extracted_title
        
    except Exception as e:
        print(f"❌ Failed to ingest {ppt_path}: {e}")
        return None


def ingest_all_ppts(demo_video_placeholder="demo://placeholder_video"):
    """
    Auto-discovers and ingests all PPT files in the data/ppts/ directory.
    """
    if not os.path.exists(PPT_DIR):
        print(f"❌ Error: {PPT_DIR} directory not found")
        return
    
    ppt_files = [f for f in os.listdir(PPT_DIR) if f.lower().endswith('.pptx')]
    
    if not ppt_files:
        print(f"⚠️  No PPT files found in {PPT_DIR}")
        return
    
    print(f"📂 Found {len(ppt_files)} PPT file(s)\n")
    
    for ppt_file in ppt_files:
        ppt_path = os.path.join(PPT_DIR, ppt_file)
        ingest_ppt(ppt_path, demo_video_placeholder)
    
    print(f"\n📄 Manifest saved to: {MANIFEST_PATH}")


if __name__ == "__main__":
    ingest_all_ppts()