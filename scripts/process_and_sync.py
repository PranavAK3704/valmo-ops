"""
process_and_sync.py
====================
Runs inside GitHub Actions.
Reads environment variables set as GitHub Secrets.

Flow:
1. Fetch rows from Training_Input sheet
2. Download each PPT from Drive
3. Run extraction (identical logic to ppt_process_map_builder.py)
4. Write results to Training_Videos sheet
5. Extension picks up changes via sheets-sync.js

All extraction functions are IDENTICAL to ppt_process_map_builder.py.
"""

import os
import re
import csv
import json
import sys
import tempfile
import traceback
from io import StringIO

import requests
import gspread
from google.oauth2.service_account import Credentials
from pptx import Presentation

# ─── Add scripts/ to path so we can import existing modules ───
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), 'scripts'))
from ppt_segmenter import segment_processes
from platform_detector import detect_platform

# ═══════════════════════════════════════════════════════════════
# CONFIG - from GitHub Secrets (set as env vars in workflow)
# ═══════════════════════════════════════════════════════════════

SHEET_ID            = os.environ["SHEET_ID"]
SHEET_INPUT_CSV_URL = os.environ["SHEET_INPUT_CSV_URL"]
SA_JSON             = os.environ["GOOGLE_SERVICE_ACCOUNT_JSON"]

INPUT_TAB_NAME  = "Training_Input"
OUTPUT_TAB_NAME = "Training_Videos"

# ═══════════════════════════════════════════════════════════════
# GOOGLE SHEETS WRITE-BACK
# ═══════════════════════════════════════════════════════════════

def get_sheets_client():
    creds_dict = json.loads(SA_JSON)
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive.readonly"
    ]
    creds = Credentials.from_service_account_info(creds_dict, scopes=scopes)
    return gspread.authorize(creds)


def write_results_to_sheet(all_log10):
    """
    Writes extracted processes to Training_Videos tab.
    Clears and rewrites completely each run (idempotent).
    """
    print(f"\n📝 Writing {len(all_log10)} processes to '{OUTPUT_TAB_NAME}' tab...")

    client    = get_sheets_client()
    sheet     = client.open_by_key(SHEET_ID)

    try:
        worksheet = sheet.worksheet(OUTPUT_TAB_NAME)
    except gspread.WorksheetNotFound:
        worksheet = sheet.add_worksheet(title=OUTPUT_TAB_NAME, rows=500, cols=10)

    headers = ["Process_Name", "URL_Module", "Start_Tab", "Video_Link", "Platform", "Active"]

    rows = [headers]
    for proc in all_log10:
        rows.append([
            proc.get("process_name", ""),
            proc.get("url_module", "") or "",
            proc.get("start_tab", ""),
            proc.get("video_link", ""),
            "log10",
            "TRUE"
        ])

    worksheet.clear()
    worksheet.update("A1", rows)

    print(f"✅ Written {len(all_log10)} processes to '{OUTPUT_TAB_NAME}'")


def fetch_input_rows():
    """Fetches rows from published Training_Input CSV."""
    print("🌐 Fetching input rows from Google Sheets...")
    response = requests.get(SHEET_INPUT_CSV_URL, timeout=15)
    response.raise_for_status()
    reader = csv.DictReader(StringIO(response.text))
    rows   = list(reader)
    print(f"✅ Retrieved {len(rows)} rows")
    return rows


# ═══════════════════════════════════════════════════════════════
# ALL EXTRACTION LOGIC - IDENTICAL TO ppt_process_map_builder.py
# ═══════════════════════════════════════════════════════════════

RED_THRESHOLD = (150, 0, 0)


def load_tab_url_map():
    """
    Loads tab->module mapping from log10_tab_url_map.csv.
    Path is resolved relative to repo root so it always works in Actions.
    """
    repo_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    csv_path  = os.path.join(repo_root, "data", "log10_tab_url_map.csv")

    mapping = {}

    if not os.path.exists(csv_path):
        print(f"⚠️  Tab-URL map not found at {csv_path} — url_module will be None")
        return mapping

    with open(csv_path, newline="", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            tab    = row.get("tab", "").strip().lower()
            url    = row.get("url", "").strip()
            module = _extract_module_from_url(url)
            if tab and module:
                mapping[tab] = module

    print(f"📋 Loaded {len(mapping)} tab→module mappings")
    return mapping


def _extract_module_from_url(url):
    try:
        path  = url.split("?")[0]
        parts = [p for p in path.split("/") if p]
        idx   = parts.index("appv2")
        return parts[idx + 1]
    except (ValueError, IndexError):
        return None


def resolve_url_module(start_tab, tab_map):
    """IDENTICAL to ppt_process_map_builder.py"""
    needle = start_tab.strip().lower()

    if needle in tab_map:
        return tab_map[needle]

    candidates = [(k, v) for k, v in tab_map.items() if needle in k]
    if candidates:
        candidates.sort(key=lambda x: len(x[0]), reverse=True)
        return candidates[0][1]

    candidates = [(k, v) for k, v in tab_map.items() if k in needle]
    if candidates:
        candidates.sort(key=lambda x: len(x[0]), reverse=True)
        return candidates[0][1]

    return None


def extract_drive_file_id(url):
    """IDENTICAL to ppt_process_map_builder.py"""
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    match = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    raise ValueError(f"Could not extract file ID from: {url}")


def download_drive_file(drive_url, save_path):
    """IDENTICAL to ppt_process_map_builder.py"""
    file_id      = extract_drive_file_id(drive_url)
    download_url = f"https://drive.google.com/uc?export=download&id={file_id}"

    session  = requests.Session()
    response = session.get(download_url, stream=True)

    for key, value in response.cookies.items():
        if key.startswith('download_warning'):
            download_url = f"https://drive.google.com/uc?export=download&id={file_id}&confirm={value}"
            response     = session.get(download_url, stream=True)
            break

    with open(save_path, "wb") as f:
        for chunk in response.iter_content(chunk_size=32768):
            if chunk:
                f.write(chunk)

    file_size = os.path.getsize(save_path)
    if file_size < 1000:
        with open(save_path, 'rb') as f:
            header = f.read(100)
            if b'<!DOCTYPE' in header or b'<html' in header:
                raise ValueError("Downloaded HTML instead of PPT - check Drive permissions")

    print(f"⬇️  Downloaded PPT ({file_size:,} bytes)")


def clean_process_name(title_text):
    """IDENTICAL to ppt_process_map_builder.py"""
    name = title_text.replace('\n', ' ').replace('\r', ' ')
    name = ' '.join(name.split())
    for suffix in ['Training Material', 'Training', 'Process', 'SOP', 'Procedure']:
        if name.endswith(suffix):
            name = name[:-len(suffix)].strip()
    return name


def clean_tab_name(raw_tab):
    """IDENTICAL to ppt_process_map_builder.py"""
    if not raw_tab:
        return None
    tab = raw_tab.replace('\n', ' ').replace('\r', ' ')
    tab = ' '.join(tab.split())
    tab = re.sub(r'^\d+\.\s*', '', tab)
    for verb in ['Navigate to', 'navigate to', 'Go to', 'go to',
                 'Click on', 'click on', 'Then click', 'then click',
                 'Click', 'click', 'Select', 'select', 'Open', 'open',
                 'Choose', 'choose']:
        if tab.startswith(verb):
            tab = tab[len(verb):].strip()
            break
    tab = re.sub(r'\s*\(\d+\)\s*$', '', tab)
    return tab.strip()


def is_red_outlined_shape(shape):
    """IDENTICAL to ppt_process_map_builder.py"""
    try:
        if not hasattr(shape, 'line'):
            return False
        if not shape.line.color.rgb:
            return False
        r, g, b = shape.line.color.rgb
        return r >= RED_THRESHOLD[0] and g < 100 and b < 100
    except:
        return False


def extract_text_from_shape(shape):
    """IDENTICAL to ppt_process_map_builder.py"""
    try:
        if hasattr(shape, 'text') and shape.text:
            return shape.text.strip()
    except:
        pass
    return None


def find_nearest_text(slide, ref_shape, max_distance=2000000):
    """IDENTICAL to ppt_process_map_builder.py"""
    min_dist  = float("inf")
    best_text = None
    for shape in slide.shapes:
        if shape == ref_shape:
            continue
        text = extract_text_from_shape(shape)
        if not text or len(text) > 100:
            continue
        dist = abs(shape.left - ref_shape.left) + abs(shape.top - ref_shape.top)
        if dist < min_dist and dist < max_distance:
            min_dist  = dist
            best_text = text
    return best_text


def extract_instruction_text(slide):
    """IDENTICAL to ppt_process_map_builder.py"""
    for shape in slide.shapes:
        text = extract_text_from_shape(shape)
        if text and len(text) > 50:
            if any(m in text.lower() for m in ['go to', 'click on', 'click', 'then', 'select']):
                return text
    return None


def parse_instruction_sequence(instruction_text):
    """IDENTICAL to ppt_process_map_builder.py"""
    if not instruction_text:
        return []
    matches     = re.findall(r'([A-Za-z][A-Za-z\s]+)\(\d+\)', instruction_text)
    clean_steps = []
    for match in matches:
        cleaned = clean_tab_name(match)
        if cleaned and cleaned not in clean_steps:
            clean_steps.append(cleaned)
    return clean_steps


def extract_log10_steps(segment, prs):
    """IDENTICAL to ppt_process_map_builder.py"""
    steps                = []
    start_tab            = None
    instruction_sequence = []

    for slide_idx in segment["Slides"]:
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]

        instruction_text = extract_instruction_text(slide)
        if instruction_text:
            sequence = parse_instruction_sequence(instruction_text)
            if sequence:
                instruction_sequence.extend(sequence)

        for shape in slide.shapes:
            if is_red_outlined_shape(shape):
                text = find_nearest_text(slide, shape)
                if text:
                    cleaned = clean_tab_name(text)
                    if cleaned and cleaned not in steps:
                        steps.append(cleaned)

    if instruction_sequence:
        start_tab = instruction_sequence[0]
        steps     = instruction_sequence
    elif steps:
        start_tab = steps[0]

    return start_tab, steps


def build_process_map(ppt_path, demo_video_link):
    """IDENTICAL to ppt_process_map_builder.py"""
    prs      = Presentation(ppt_path)
    segments = segment_processes(ppt_path)
    tab_map  = load_tab_url_map()

    if not segments:
        print("⚠️  No processes detected in PPT")
        return {"log10": [], "external": []}

    log10_processes    = []
    external_processes = []

    for segment in segments:
        platform     = detect_platform(segment, ppt_path)
        process_name = clean_process_name(segment["Process_Name"])

        print(f"\n🔍 {process_name} | Platform: {platform}")

        if platform == "log10":
            try:
                start_tab, steps = extract_log10_steps(segment, prs)

                if start_tab:
                    url_module = resolve_url_module(start_tab, tab_map)
                    log10_processes.append({
                        "process_name": process_name,
                        "platform":     "log10",
                        "start_tab":    start_tab,
                        "url_module":   url_module,
                        "steps":        steps,
                        "video_link":   demo_video_link,
                        **({"needs_review": True} if not url_module else {})
                    })
                    print(f"   ✅ Start Tab: {start_tab} | URL Module: {url_module}")
                else:
                    log10_processes.append({
                        "process_name": process_name,
                        "platform":     "log10",
                        "start_tab":    "Dashboard",
                        "url_module":   resolve_url_module("Dashboard", tab_map),
                        "steps":        [],
                        "video_link":   demo_video_link,
                        "needs_review": True
                    })
                    print(f"   ⚠️  Could not extract start tab - needs review")

            except Exception as e:
                print(f"   ❌ Error: {e}")
                log10_processes.append({
                    "process_name": process_name,
                    "platform":     "log10",
                    "start_tab":    "Dashboard",
                    "url_module":   resolve_url_module("Dashboard", tab_map),
                    "steps":        [],
                    "video_link":   demo_video_link,
                    "needs_review": True
                })
        else:
            external_processes.append({
                "process_name": process_name,
                "platform":     platform if platform != "unknown" else "external",
                "video_link":   demo_video_link,
                "use_case":     "training_only"
            })
            print(f"   ℹ️  External process - skipping")

    return {"log10": log10_processes, "external": external_processes}


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    print("\n" + "="*60)
    print("🚀 Valmo PPT Processor - GitHub Actions")
    print("="*60)

    all_log10    = []
    all_external = []
    errors       = []

    rows = fetch_input_rows()

    if not rows:
        print("❌ No rows in Training_Input sheet")
        sys.exit(1)

    for idx, row in enumerate(rows, start=1):
        process_name = row.get("process_name", "").strip()
        ppt_link     = row.get("ppt_link", "").strip()
        video_link   = row.get("video_link", "").strip()

        if not ppt_link:
            print(f"\n⚠️  Row {idx}: No PPT link, skipping")
            continue

        if not process_name:
            process_name = f"Process_{idx}"

        print(f"\n{'='*60}")
        print(f"Row {idx}: {process_name}")
        print(f"{'='*60}")

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            temp_path = tmp.name

        try:
            download_drive_file(ppt_link, temp_path)
            results = build_process_map(temp_path, video_link)
            all_log10.extend(results["log10"])
            all_external.extend(results["external"])

        except Exception as e:
            print(f"❌ Error on row {idx}: {e}")
            traceback.print_exc()
            errors.append({"row": idx, "process": process_name, "error": str(e)})

        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    # Write to sheet
    write_results_to_sheet(all_log10)

    # Summary
    print(f"\n{'='*60}")
    print(f"✅ DONE")
    print(f"{'='*60}")
    print(f"Log10 processes:    {len(all_log10)}")
    print(f"External processes: {len(all_external)}")

    if all_log10:
        print(f"\n📊 Log10 Processes:")
        for proc in all_log10:
            flag = " ⚠️ NEEDS REVIEW" if proc.get("needs_review") else ""
            print(f"  • {proc['process_name']}")
            print(f"    Tab: {proc['start_tab']} | Module: {proc.get('url_module', '⚠️ None')}{flag}")

    if errors:
        print(f"\n⚠️  {len(errors)} errors:")
        for e in errors:
            print(f"  Row {e['row']} ({e['process']}): {e['error']}")
        sys.exit(1)  # Fail the action so you see it in GitHub


if __name__ == "__main__":
    main()
